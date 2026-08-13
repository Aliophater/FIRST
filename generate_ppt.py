# -*- coding: utf-8 -*-
"""
速云记 — 笔记智能体助手 答辩PPT生成脚本
项目工期: 2026.06.28 - 2026.07.13
答辩日期: 2026.07.16
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
import copy

# ======================== 颜色定义 ========================
PRIMARY = RGBColor(0x1A, 0x3C, 0x6E)       # 深蓝主色
ACCENT = RGBColor(0x2E, 0x86, 0xC1)        # 亮蓝强调
ACCENT2 = RGBColor(0x16, 0xA0, 0x85)       # 青绿
ACCENT3 = RGBColor(0xE8, 0x6A, 0x33)       # 橙色
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF0, 0xF4, 0xF8)
DARK_TEXT = RGBColor(0x2C, 0x3E, 0x50)
GRAY = RGBColor(0x7F, 0x8C, 0x8D)
LIGHT_GRAY = RGBColor(0xEC, 0xF0, 0xF1)
CARD_BG = RGBColor(0xE8, 0xEE, 0xF2)
TABLE_HEADER = RGBColor(0x1A, 0x3C, 0x6E)
TABLE_ROW1 = RGBColor(0xF8, 0xFA, 0xFC)
TABLE_ROW2 = RGBColor(0xEC, 0xF0, 0xF4)

# ======================== 创建演示文稿 ========================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

# ======================== 辅助函数 ========================

def add_blank_slide():
    """添加空白幻灯片"""
    layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(layout)
    return slide

def set_bg_color(slide, color):
    """设置幻灯片背景色"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, line_color=None, shape_type=MSO_SHAPE.RECTANGLE):
    """添加形状"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=14, color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multi_text(slide, left, top, width, height, lines, font_size=14, color=DARK_TEXT, font_name='Microsoft YaHei', line_spacing=1.5):
    """添加多行文本"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(font_size * 0.6)
        p.space_before = Pt(font_size * 0.2)
    return txBox

def add_title_bar(slide, title_text, subtitle_text=None):
    """添加页面标题栏"""
    # 顶部色条
    bar = add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.9), PRIMARY)
    # 左侧装饰条
    deco = add_shape(slide, Inches(0), Inches(0), Inches(0.12), Inches(0.9), ACCENT2)
    # 标题文字
    add_textbox(slide, Inches(0.4), Inches(0.12), Inches(8), Inches(0.5), title_text, font_size=26, color=WHITE, bold=True)
    if subtitle_text:
        add_textbox(slide, Inches(0.4), Inches(0.55), Inches(8), Inches(0.35), subtitle_text, font_size=12, color=RGBColor(0xBD, 0xD7, 0xEE))
    # 右侧页码占位
    add_textbox(slide, Inches(12), Inches(0.25), Inches(1.2), Inches(0.4), "", font_size=11, color=RGBColor(0xBD, 0xD7, 0xEE), alignment=PP_ALIGN.RIGHT)

def add_card(slide, left, top, width, height, title, content_lines, title_color=ACCENT, content_size=11, title_size=13):
    """添加卡片式内容块"""
    # 卡片背景
    card = add_shape(slide, left, top, width, height, CARD_BG)
    card.shadow.inherit = False
    # 顶部色条
    top_bar = add_shape(slide, left, top, width, Inches(0.04), title_color)
    # 标题
    add_textbox(slide, left + Inches(0.15), top + Inches(0.1), width - Inches(0.3), Inches(0.35), title, font_size=title_size, color=title_color, bold=True)
    # 内容
    if content_lines:
        add_multi_text(slide, left + Inches(0.15), top + Inches(0.5), width - Inches(0.3), height - Inches(0.6), content_lines, font_size=content_size, color=DARK_TEXT)
    return card

def add_table(slide, left, top, width, rows_data, col_widths=None, header=True, font_size=11):
    """添加表格"""
    num_rows = len(rows_data)
    num_cols = len(rows_data[0]) if rows_data else 0
    height = Inches(0.4) * num_rows

    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for row_idx, row_data in enumerate(rows_data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_text)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.name = 'Microsoft YaHei'
                if row_idx == 0 and header:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                else:
                    p.font.color.rgb = DARK_TEXT
                p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if row_idx == 0 and header:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ROW1 if row_idx % 2 == 1 else TABLE_ROW2

    return table_shape

def add_arrow_connector(slide, start_x, start_y, end_x, end_y, color=ACCENT):
    """添加箭头连接线"""
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, start_x, start_y, end_x, end_y)
    connector.line.color.rgb = color
    connector.line.width = Pt(2)
    # 设置箭头
    line_elem = connector.line._get_or_add_ln()
    tail_end = line_elem.makeelement(qn('a:tailEnd'), {})
    tail_end.set('type', 'triangle')
    tail_end.set('w', 'med')
    tail_end.set('h', 'med')
    line_elem.append(tail_end)
    return connector

def add_rounded_box(slide, left, top, width, height, text, fill_color, text_color=WHITE, font_size=12, bold=True):
    """添加圆角矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = bold
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    shape.text_frame.margin_top = Pt(0)
    shape.text_frame.margin_bottom = Pt(0)
    return shape

def add_page_number(slide, num):
    """添加页码"""
    add_textbox(slide, Inches(12.3), Inches(7.1), Inches(0.9), Inches(0.3), f"{num}", font_size=10, color=GRAY, alignment=PP_ALIGN.RIGHT)

# ======================== 开始生成PPT ========================

# ============================================================
# 第1页: 封面
# ============================================================
slide = add_blank_slide()
set_bg_color(slide, PRIMARY)

# 装饰几何图形
deco1 = add_shape(slide, Inches(-1), Inches(-1), Inches(5), Inches(5), ACCENT, shape_type=MSO_SHAPE.OVAL)
deco1.fill.fore_color.rgb = RGBColor(0x22, 0x50, 0x8A)
deco2 = add_shape(slide, Inches(10), Inches(4), Inches(6), Inches(6), ACCENT2, shape_type=MSO_SHAPE.OVAL)
deco2.fill.fore_color.rgb = RGBColor(0x16, 0x80, 0x6E)

# 项目标签
tag = add_rounded_box(slide, Inches(4.6), Inches(1.2), Inches(4.1), Inches(0.45), "AI驱动的个人知识管理工具", ACCENT2, WHITE, 13)

# 主标题
add_textbox(slide, Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.0), "速云记", font_size=52, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(3.0), Inches(10.3), Inches(0.7), "笔记智能体助手 — 项目答辩报告", font_size=24, color=RGBColor(0xBD, 0xD7, 0xEE), alignment=PP_ALIGN.CENTER)

# 分割线
line = add_shape(slide, Inches(4.5), Inches(3.9), Inches(4.3), Inches(0.03), ACCENT)

# 工期与答辩信息
add_textbox(slide, Inches(2), Inches(4.2), Inches(9.3), Inches(0.4), "项目工期: 2026年6月28日 — 2026年7月13日  |  答辩日期: 2026年7月16日", font_size=14, color=RGBColor(0x9B, 0xC4, 0xDA), alignment=PP_ALIGN.CENTER)

# 团队信息
add_textbox(slide, Inches(2), Inches(5.0), Inches(9.3), Inches(0.4), "团队成员: 沈骏龙  余致强  李子恩  贠炳鑫  闫浩楠  崔浩文", font_size=15, color=WHITE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(2), Inches(5.5), Inches(9.3), Inches(0.4), "指导团队: 西北大学  软件工程专业", font_size=12, color=RGBColor(0x9B, 0xC4, 0xDA), alignment=PP_ALIGN.CENTER)

# 技术栈标签
techs = ["FastAPI", "LangChain", "React 19", "ChromaDB", "MySQL", "Redis"]
for i, tech in enumerate(techs):
    x = Inches(3.0 + i * 1.3)
    add_rounded_box(slide, x, Inches(6.2), Inches(1.15), Inches(0.35), tech, ACCENT, WHITE, 10)

# ============================================================
# 第2页: 目录
# ============================================================
slide = add_blank_slide()
set_bg_color(slide, LIGHT_BG)
add_title_bar(slide, "目  录", "CONTENTS")

contents = [
    ("01", "项目概述", "Project Overview"),
    ("02", "成员及分工", "Team & Division"),
    ("03", "技术架构与选型", "Tech Architecture"),
    ("04", "数据库设计", "Database Design"),
    ("05", "项目功能设计", "Functional Design"),
    ("06", "核心功能实现", "Core Implementation"),
    ("07", "系统运行展示", "System Demo"),
    ("08", "遇到的问题", "Challenges"),
    ("09", "总结与展望", "Summary & Outlook"),
]

for i, (num, title, en) in enumerate(contents):
    row = i // 3
    col = i % 3
    x = Inches(0.8 + col * 4.2)
    y = Inches(1.3 + row * 1.8)

    # 数字圆圈
    circle = add_shape(slide, x, y, Inches(0.8), Inches(0.8), ACCENT, shape_type=MSO_SHAPE.OVAL)
    tf = circle.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER
    circle.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 标题
    add_textbox(slide, x + Inches(1.0), y + Inches(0.05), Inches(3), Inches(0.4), title, font_size=16, color=PRIMARY, bold=True)
    add_textbox(slide, x + Inches(1.0), y + Inches(0.45), Inches(3), Inches(0.3), en, font_size=10, color=GRAY)

add_page_number(slide, 2)

# ============================================================
# 第3页: 一、项目概述
# ============================================================
slide = add_blank_slide()
set_bg_color(slide, WHITE)
add_title_bar(slide, "一、项目概述", "Project Overview")

# 项目简介
add_textbox(slide, Inches(0.5), Inches(1.1), Inches(12), Inches(0.4), "项目简介", font_size=18, color=PRIMARY, bold=True)
intro_lines = [
    "速云记是一款基于 FastAPI + LangChain 构建的智能笔记助手，融合「笔记管理 + RAG知识库 + AI写作辅助」三大核心能力。",
    "系统旨在解决「笔记写了从不回看、知识散落成孤岛」的痛点，通过AI赋能实现笔记的智能化管理、语义检索与个性化回顾。",
]
add_multi_text(slide, Inches(0.5), Inches(1.55), Inches(12.3), Inches(0.8), intro_lines, font_size=13, color=DARK_TEXT)

# 核心能力卡片
add_textbox(slide, Inches(0.5), Inches(2.5), Inches(12), Inches(0.4), "核心能力", font_size=18, color=PRIMARY, bold=True)

cards = [
    ("笔记管理", ["Markdown编辑器", "智能标签(LLM自动分类)", "语义搜索", "Markdown导出"], ACCENT),
    ("RAG知识库", ["多格式文档上传", "txt/pdf/md/pptx/docx", "向量检索精准问答", "文档引用来源展示"], ACCENT2),
    ("AI写作辅助", ["联机补全(Tab键采纳)", "续写/扩写/摘要生成", "SSE流式输出", "跨源关联推荐"], ACCENT3),
    ("间隔重复回顾", ["艾宾浩斯遗忘曲线", "1/2/4/7/15/30天间隔", "LLM生成回顾选择题", "自动回顾提醒"], PRIMARY),
]
for i, (title, items, color) in enumerate(cards):
    x = Inches(0.5 + i * 3.15)
    add_card(slide, x, Inches(3.0), Inches(2.95), Inches(2.5), title, [f"• {item}" for item in items], title_color=color, content_size=11, title_size=14)

# 技术亮点
add_textbox(slide, Inches(0.5), Inches(5.7), Inches(12), Inches(0.4), "技术亮点", font_size=18, color=PRIMARY, bold=True)
highlights = [
    "微服务架构: FastAPI后端 + Django用户服务 + React前端，三服务独立部署",
    "Agent工具系统: 9个LangChain Tool，支持笔记CRUD、RAG检索、回顾管理等Agent自主调用",
    "HyDE混合检索: 假设性文档生成 + BM25 + 向量检索动态权重 + CrossEncoder重排序",
    "用户级数据隔离: JWT认证 + ChromaDB metadata过滤，RAG检索只能访问本人数据",
]
add_multi_text(slide, Inches(0.5), Inches(6.15), Inches(12.3), Inches(1.2), [f"▸ {h}" for h in highlights], font_size=12, color=DARK_TEXT, line_spacing=1.3)

add_page_number(slide, 3)

# ============================================================
# 第4页: 二、成员及分工 - 组织结构图
# ============================================================
slide = add_blank_slide()
set_bg_color(slide, WHITE)
add_title_bar(slide, "二、成员及分工", "Team & Division — 组织结构图")

# 组织结构图
# 组长节点
leader = add_rounded_box(slide, Inches(5.4), Inches(1.2), Inches(2.5), Inches(0.6), "组长: 沈骏龙\nAI Agent核心 + 后端架构", PRIMARY, WHITE, 12)

# 连接线 - 组长到5个组员
member_y = Inches(3.0)
positions = [
    (Inches(0.4), "余致强", "RAG知识库系统\n多格式文档解析\nChromaDB存储\nHyDE混合检索\n重排序模型", ACCENT),
    (Inches(2.95), "李子恩", "前端核心与编辑器\nReact架构\nTiptap富文本编辑器\n笔记管理页面\n通用组件库", ACCENT2),
    (Inches(5.5), "贠炳鑫", "前端AI交互与国际化\nAI聊天+SSE流式\n知识库页面\n每日回顾\ni18n/Zustand", ACCENT3),
    (Inches(8.05), "闫浩楠", "用户服务与数据库\nDjango用户系统\nJWT认证\nMySQL表设计\nRedis缓存", RGBColor(0x8E, 0x44, 0xAD)),
    (Inches(10.6), "崔浩文", "基础设施与展示\n路由守卫\n登录/注册页面\nDocker部署\n项目文档/PPT", RGBColor(0x2C, 0x3E, 0x50)),
]

for x, name, desc, color in positions:
    # 垂直连接线从组长到组员
    mid_x = x + Inches(1.1)
    add_arrow_connector(slide, Inches(6.65), Inches(1.8), mid_x, member_y, GRAY)
    # 组员节点
    add_rounded_box(slide, x, member_y, Inches(2.2), Inches(0.5), name, color, WHITE, 12)
    # 职责描述
    desc_lines = desc.split("\n")
    add_multi_text(slide, x, member_y + Inches(0.55), Inches(2.2), Inches(2.2), [f"• {l}" for l in desc_lines[1:]], font_size=10, color=DARK_TEXT, line_spacing=1.2)
    # 职责标题
    add_textbox(slide, x, member_y + Inches(0.5), Inches(2.2), Inches(0.3), desc_lines[0], font_size=10, color=color, bold=True)

# 底部统计
add_shape(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.8), LIGHT_BG)
stats_text = "团队规模: 6人  |  角色分配: 1名组长(后端架构) + 1名RAG工程师 + 2名前端工程师 + 1名数据库工程师 + 1名DevOps"
add_textbox(slide, Inches(0.8), Inches(6.45), Inches(12), Inches(0.5), stats_text, font_size=12, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

add_page_number(slide, 4)

# ============================================================
# 第5页: 三、技术架构与选型 - 系统技术架构图
# ============================================================
slide = add_blank_slide()
set_bg_color(slide, WHITE)
add_title_bar(slide, "三、技术架构与选型", "Tech Architecture — 分层架构图")

# 分层架构图
layers = [
    ("前端展示层", "React 19 + TypeScript + Vite + Tailwind CSS + Tiptap + Zustand + i18next", ACCENT3, Inches(1.1)),
    ("API网关层", "FastAPI (端口8000)  |  CORS中间件 + JWT鉴权 + 限流(令牌桶) + 统一响应封装", ACCENT, Inches(2.1)),
    ("业务服务层", "笔记服务  |  RAG服务  |  Agent工厂  |  回顾服务  |  知识库服务  |  会话管理", ACCENT2, Inches(3.1)),
    ("AI/LLM层", "LangChain Agent (9 Tools)  |  Qwen3-Max (阿里云百炼)  |  Ollama qwen3.5:0.8b  |  bge-reranker-v2-m3", RGBColor(0x8E, 0x44, 0xAD), Inches(4.1)),
    ("数据存储层", "MySQL 8.0 (笔记/回顾/会话)  |  ChromaDB (rag_collection + notes_collection)  |  Redis 7 (缓存)", PRIMARY, Inches(5.1)),
    ("基础设施层", "Docker Compose  |  Django用户服务 (端口8001)  |  Celery异步任务  |  LangSmith追踪", RGBColor(0x2C, 0x3E, 0x50), Inches(6.1)),
]

for name, desc, color, y in layers:
    # 层标签
    add_rounded_box(slide, Inches(0.5), y, Inches(2.5), Inches(0.8), name, color, WHITE, 13)
    # 层内容
    box = add_shape(slide, Inches(3.2), y, Inches(9.6), Inches(0.8), LIGHT_GRAY)
    box.line.color.rgb = color
    box.line.width = Pt(1)
    box.shadow.inherit = False
    add_textbox(slide, Inches(3.4), y + Inches(0.15), Inches(9.2), Inches(0.5), desc, font_size=11, color=DARK_TEXT)

# 右侧技术选型表
add_page_number(slide, 5)

# ============================================================
# 第6页: 技术选型详细表
# ============================================================
slide = add_blank_slide()
set_bg_color(slide, WHITE)
add_title_bar(slide, "三、技术架构与选型", "Tech Stack — 技术选型清单")

# 后端技术
add_textbox(slide, Inches(0.4), Inches(1.1), Inches(6), Inches(0.35), "后端技术栈", font_size=15, color=PRIMARY, bold=True)
backend_rows = [
    ["技术", "版本", "用途"],
    ["FastAPI", "0.115+", "高性能异步Web框架，API路由+CORS+中间件"],
    ["LangChain", "0.3+", "Agent应用开发框架（AgentExecutor + Tools）"],
    ["SQLAlchemy", "2.0+", "异步ORM，管理MySQL（aiomysql驱动）"],
    ["ChromaDB", "0.5+", "轻量级向量数据库（rag_collection + notes_collection）"],
    ["DashScope API", "Qwen3-Max", "阿里云百炼大模型服务（Agent + 摘要）"],
    ["Ollama", "qwen3.5:0.8b", "本地模型部署（联机补全，低延迟）"],
    ["bge-reranker-v2-m3", "CrossEncoder", "HuggingFace重排序模型"],
    ["Redis", "7.0+", "缓存 + Celery消息队列"],
]
add_table(slide, Inches(0.4), Inches(1.5), Inches(6.2), backend_rows, col_widths=[Inches(1.6), Inches(1.3), Inches(3.3)], font_size=10)

# 前端技术
add_textbox(slide, Inches(6.9), Inches(1.1), Inches(6), Inches(0.35), "前端技术栈", font_size=15, color=PRIMARY, bold=True)
frontend_rows = [
    ["技术", "版本", "用途"],
    ["React", "19", "现代化前端框架"],
    ["TypeScript", "5.x", "类型安全"],
    ["Vite", "6.x", "极速构建工具"],
    ["Tailwind CSS", "4.x", "原子化CSS框架"],
    ["Tiptap", "2.x", "富文本Markdown编辑器"],
    ["Radix UI", "latest", "无头UI组件库"],
    ["Zustand", "5.x", "轻量状态管理"],
    ["i18next", "24.x", "国际化（中/英）"],
    ["Axios", "1.x", "HTTP客户端"],
]
add_table(slide, Inches(6.9), Inches(1.5), Inches(5.9), frontend_rows, col_widths=[Inches(1.4), Inches(1.2), Inches(3.3)], font_size=10)

# 选型理由
add_textbox(slide, Inches(0.4), Inches(5.5), Inches(12), Inches(0.35), "关键选型理由", font_size=15, color=PRIMARY, bold=True)
reasons = [
    "FastAPI vs Django: 选择FastAPI作为主服务因其原生异步支持，适合Agent流式输出和RAG高并发检索场景；Django仅用于用户认证微服务",
    "ChromaDB vs Milvus: ChromaDB轻量嵌入式，无需额外部署服务，适合项目规模；持久化到本地磁盘，开发调试便捷",
    "双LLM策略: 阿里云Qwen3-Max用于Agent推理和摘要生成（高质量），本地Ollama用于联机补全（低延迟300-500ms）",
    "HyDE + 混合检索: HyDE提升语义匹配准确率，BM25+向量动态权重兼顾关键词和语义，CrossEncoder重排序精排Top-3",
]
add_multi_text(slide, Inches(0.4), Inches(5.9), Inches(12.5), Inches(1.5), [f"▸ {r}" for r in reasons], font_size=11, color=DARK_TEXT, line_spacing=1.3)

add_page_number(slide, 6)

# ============================================================
# 第7页: 四、数据库设计 - ER图
# ============================================================
slide = add_blank_slide()
set_bg_color(slide, WHITE)
add_title_bar(slide, "四、数据库设计", "Database Design — ER图")

# ER图 - 5个实体
# User实体 (Django user_service DB)
add_rounded_box(slide, Inches(0.5), Inches(1.2), Inches(3.2), Inches(2.8), "", RGBColor(0xF0, 0xF4, 0xF8), PRIMARY, 12)
add_shape(slide, Inches(0.5), Inches(1.2), Inches(3.2), Inches(0.4), RGBColor(0x8E, 0x44, 0xAD))
add_textbox(slide, Inches(0.5), Inches(1.23), Inches(3.2), Inches(0.35), "User (用户服务DB)", font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
user_fields = [
    "uuid (PK, ShortUUID)",
    "username (varchar 150)",
    "email (unique)",
    "telephone (varchar 11)",
    "status (0/1/2)",
    "gender (1/2/3)",
    "bio (text)",
    "avatar (varchar 255)",
    "date_joined (datetime)",
    "last_login (datetime)",
]
add_multi_text(slide, Inches(0.65), Inches(1.7), Inches(2.9), Inches(2.2), user_fields, font_size=9, color=DARK_TEXT, line_spacing=1.1)

# Note实体 (chat_history DB)
add_rounded_box(slide, Inches(5.0), Inches(1.2), Inches(3.2), Inches(2.5), "", RGBColor(0xF0, 0xF4, 0xF8), PRIMARY, 12)
add_shape(slide, Inches(5.0), Inches(1.2), Inches(3.2), Inches(0.4), ACCENT)
add_textbox(slide, Inches(5.0), Inches(1.23), Inches(3.2), Inches(0.35), "Note (笔记表)", font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
note_fields = [
    "id (PK, UUID)",
    "user_id (FK → User.uuid)",
    "title (varchar 200)",
    "content (text, Markdown)",
    "tags (JSON数组)",
    "category (work/study/...)",
    "created_at (datetime)",
    "updated_at (datetime)",
]
add_multi_text(slide, Inches(5.15), Inches(1.7), Inches(2.9), Inches(1.9), note_fields, font_size=9, color=DARK_TEXT, line_spacing=1.1)

# ReviewRecord实体
add_rounded_box(slide, Inches(9.5), Inches(1.2), Inches(3.3), Inches(2.3), "", RGBColor(0xF0, 0xF4, 0xF8), PRIMARY, 12)
add_shape(slide, Inches(9.5), Inches(1.2), Inches(3.3), Inches(0.4), ACCENT2)
add_textbox(slide, Inches(9.5), Inches(1.23), Inches(3.3), Inches(0.35), "ReviewRecord (回顾记录)", font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
review_fields = [
    "id (PK, UUID)",
    "note_id (FK → Note.id CASCADE)",
    "user_id (FK → User.uuid)",
    "review_count (int)",
    "next_review_at (datetime)",
    "last_reviewed_at (datetime)",
    "interval_days (int)",
    "created_at (datetime)",
]
add_multi_text(slide, Inches(9.65), Inches(1.7), Inches(3.0), Inches(1.7), review_fields, font_size=9, color=DARK_TEXT, line_spacing=1.1)

# ChatSession实体
add_rounded_box(slide, Inches(2.0), Inches(4.3), Inches(3.5), Inches(2.3), "", RGBColor(0xF0, 0xF4, 0xF8), PRIMARY, 12)
add_shape(slide, Inches(2.0), Inches(4.3), Inches(3.5), Inches(0.4), ACCENT3)
add_textbox(slide, Inches(2.0), Inches(4.33), Inches(3.5), Inches(0.35), "ChatSession (会话表)", font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
session_fields = [
    "id (PK, varchar 64)",
    "user_id (FK → User.uuid)",
    "title (varchar 255)",
    "metadata (JSON)",
    "created_at (datetime)",
    "updated_at (datetime)",
]
add_multi_text(slide, Inches(2.15), Inches(4.8), Inches(3.2), Inches(1.7), session_fields, font_size=9, color=DARK_TEXT, line_spacing=1.1)

# ChatMessage实体
add_rounded_box(slide, Inches(7.5), Inches(4.3), Inches(3.5), Inches(2.5), "", RGBColor(0xF0, 0xF4, 0xF8), PRIMARY, 12)
add_shape(slide, Inches(7.5), Inches(4.3), Inches(3.5), Inches(0.4), RGBColor(0xE7, 0x4C, 0x3C))
add_textbox(slide, Inches(7.5), Inches(4.33), Inches(3.5), Inches(0.35), "ChatMessage (消息表)", font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
msg_fields = [
    "id (PK, auto increment)",
    "session_id (FK → ChatSession.id)",
    "role (user/assistant)",
    "content (text)",
    "metadata (JSON)",
    "created_at (datetime)",
    "relationship: session.messages",
]
add_multi_text(slide, Inches(7.65), Inches(4.8), Inches(3.2), Inches(1.9), msg_fields, font_size=9, color=DARK_TEXT, line_spacing=1.1)

# 关系连线
# User → Note (1:N)
add_arrow_connector(slide, Inches(3.7), Inches(2.3), Inches(5.0), Inches(2.3), PRIMARY)
add_textbox(slide, Inches(3.8), Inches(1.95), Inches(1.2), Inches(0.25), "1 : N", font_size=9, color=PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)
# Note → ReviewRecord (1:N)
add_arrow_connector(slide, Inches(8.2), Inches(2.3), Inches(9.5), Inches(2.3), ACCENT2)
add_textbox(slide, Inches(8.3), Inches(1.95), Inches(1.2), Inches(0.25), "1 : N", font_size=9, color=ACCENT2, bold=True, alignment=PP_ALIGN.CENTER)
# User → ChatSession (1:N)
add_arrow_connector(slide, Inches(2.1), Inches(3.0), Inches(3.0), Inches(4.3), GRAY)
add_textbox(slide, Inches(1.5), Inches(3.5), Inches(1.2), Inches(0.25), "1 : N", font_size=9, color=GRAY, bold=True, alignment=PP_ALIGN.CENTER)
# ChatSession → ChatMessage (1:N)
add_arrow_connector(slide, Inches(5.5), Inches(5.2), Inches(7.5), Inches(5.2), RGBColor(0xE7, 0x4C, 0x3C))
add_textbox(slide, Inches(6.0), Inches(4.85), Inches(1.2), Inches(0.25), "1 : N (CASCADE)", font_size=9, color=RGBColor(0xE7, 0x4C, 0x3C), bold=True, alignment=PP_ALIGN.CENTER)

add_page_number(slide, 7)

# ============================================================
# 第8页: 数据表清单
# ============================================================
slide = add_blank_slide()
set_bg_color(slide, WHITE)
add_title_bar(slide, "四、数据库设计", "Database Design — 数据表清单")

add_textbox(slide, Inches(0.4), Inches(1.0), Inches(12), Inches(0.3), "数据库: user_service (Django) + chat_history (FastAPI)  |  共 5 张表", font_size=12, color=GRAY)

table_rows = [
    ["表名", "所属数据库", "记录数(估计)", "主要字段", "说明"],
    ["user_service", "user_service", "—", "uuid, username, email, status, avatar", "用户信息表（Django ORM管理）"],
    ["notes", "chat_history", "—", "id, user_id, title, content, tags, category", "笔记主表（SQLAlchemy异步ORM）"],
    ["review_records", "chat_history", "—", "id, note_id, user_id, review_count, next_review_at", "间隔重复回顾记录（FK CASCADE）"],
    ["chat_sessions", "chat_history", "—", "id, user_id, title, metadata", "AI对话会话表"],
    ["chat_messages", "chat_history", "—", "id, session_id, role, content, metadata", "对话消息明细（FK CASCADE）"],
]
add_table(slide, Inches(0.4), Inches(1.5), Inches(12.5), table_rows, col_widths=[Inches(2.0), Inches(1.8), Inches(1.2), Inches(4.0), Inches(3.5)], font_size=10)

# 向量数据库说明
add_textbox(slide, Inches(0.4), Inches(4.3), Inches(12), Inches(0.35), "向量数据库: ChromaDB（双Collection设计）", font_size=15, color=PRIMARY, bold=True)
vector_rows = [
    ["Collection名称", "用途", "存储内容", "检索方式"],
    ["rag_collection", "知识库RAG检索", "上传文档的向量切片(txt/pdf/md/pptx/docx)", "HyDE + BM25混合 + 重排序"],
    ["notes_collection", "笔记语义搜索", "用户笔记的向量表示", "similarity_search + metadata过滤"],
]
add_table(slide, Inches(0.4), Inches(4.8), Inches(12.5), vector_rows, col_widths=[Inches(2.5), Inches(2.5), Inches(4.5), Inches(3.0)], font_size=10)

# 设计要点
add_textbox(slide, Inches(0.4), Inches(6.3), Inches(12), Inches(0.35), "设计要点", font_size=15, color=PRIMARY, bold=True)
points = [
    "双数据库隔离: 用户服务(user_service)与业务数据(chat_history)物理隔离，通过user_id逻辑关联，不做物理外键约束",
    "双写一致性: 笔记保存时MySQL+ChromaDB双写，更新时先删旧向量再写新向量，删除时同步清理两份数据",
    "JSON字段应用: tags存储为JSON数组(灵活扩展)，metadata存储Agent工具调用上下文，避免频繁DDL",
]
add_multi_text(slide, Inches(0.4), Inches(6.7), Inches(12.5), Inches(0.8), [f"▸ {p}" for p in points], font_size=10, color=DARK_TEXT, line_spacing=1.2)

add_page_number(slide, 8)

# ============================================================
# 第9页: 五、项目功能设计 - 功能结构图
# ============================================================
slide = add_blank_slide()
set_bg_color(slide, WHITE)
add_title_bar(slide, "五、项目功能设计", "Functional Design — 功能结构图")

# 功能结构图 - 树形
# 根节点
root = add_rounded_box(slide, Inches(5.2), Inches(1.0), Inches(3.0), Inches(0.55), "速云记 — 笔记智能体助手", PRIMARY, WHITE, 13)

# 一级模块
modules = [
    (Inches(0.3), "笔记管理", ACCENT),
    (Inches(2.85), "RAG知识库", ACCENT2),
    (Inches(5.4), "AI智能助手", ACCENT3),
    (Inches(7.95), "间隔重复回顾", RGBColor(0x8E, 0x44, 0xAD)),
    (Inches(10.5), "用户与安全", RGBColor(0x2C, 0x3E, 0x50)),
]

for x, name, color in modules:
    # 连接线
    add_arrow_connector(slide, Inches(6.7), Inches(1.55), x + Inches(1.25), Inches(2.1), GRAY)
    # 模块节点
    add_rounded_box(slide, x, Inches(2.1), Inches(2.5), Inches(0.45), name, color, WHITE, 12)

# 二级功能
sub_features = [
    # 笔记管理
    [(Inches(0.3), "新建/编辑笔记"), (Inches(0.3), "分页列表+分类筛选"), (Inches(0.3), "语义搜索"), (Inches(0.3), "批量操作"), (Inches(0.3), "Markdown导出")],
    # RAG知识库
    [(Inches(2.85), "文档上传(5种格式)"), (Inches(2.85), "切片详情查看"), (Inches(2.85), "文档管理"), (Inches(2.85), "MD5去重"), (Inches(2.85), "PDF图片提取")],
    # AI智能助手
    [(Inches(5.4), "Agent智能问答"), (Inches(5.4), "AI联机补全"), (Inches(5.4), "写作辅助(续写/扩写/摘要)"), (Inches(5.4), "跨源关联推荐"), (Inches(5.4), "SSE流式输出")],
    # 间隔重复回顾
    [(Inches(7.95), "今日回顾列表"), (Inches(7.95), "标记已回顾"), (Inches(7.95), "LLM生成选择题"), (Inches(7.95), "艾宾浩斯算法"), (Inches(7.95), "自动回顾提醒")],
    # 用户与安全
    [(Inches(10.5), "注册/登录"), (Inches(10.5), "JWT认证"), (Inches(10.5), "用户级数据隔离"), (Inches(10.5), "会话持久化"), (Inches(10.5), "个人资料管理")],
]

for module_idx, features in enumerate(sub_features):
    x_base = features[0][0]
    color = modules[module_idx][2]
    for i, (_, feat) in enumerate(features):
        y = Inches(2.75 + i * 0.75)
        # 连接线
        add_arrow_connector(slide, x_base + Inches(1.25), Inches(2.55), x_base + Inches(1.25), y, GRAY)
        # 功能项
        box = add_shape(slide, x_base, y, Inches(2.5), Inches(0.55), LIGHT_GRAY)
        box.line.color.rgb = color
        box.line.width = Pt(1)
        box.shadow.inherit = False
        add_textbox(slide, x_base + Inches(0.1), y + Inches(0.08), Inches(2.3), Inches(0.4), feat, font_size=10, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

add_page_number(slide, 9)

# ============================================================
# 第10页: 用例图
# ============================================================
slide = add_blank_slide()
set_bg_color(slide, WHITE)
add_title_bar(slide, "五、项目功能设计", "Functional Design — 用例图")

# 用例图
# 系统边界
boundary = add_shape(slide, Inches(3.0), Inches(1.0), Inches(8.5), Inches(5.8), RGBColor(0xFA, 0xFB, 0xFC), shape_type=MSO_SHAPE.RECTANGLE)
boundary.line.color.rgb = ACCENT
boundary.line.width = Pt(2)
boundary.line.dash_style = 2  # 虚线
boundary.shadow.inherit = False
add_textbox(slide, Inches(3.0), Inches(1.05), Inches(8.5), Inches(0.3), "速云记系统", font_size=12, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

# Actor - 普通用户
actor = add_shape(slide, Inches(0.8), Inches(3.0), Inches(1.5), Inches(2.0), WHITE, shape_type=MSO_SHAPE.OVAL)
actor.line.color.rgb = PRIMARY
actor.line.width = Pt(2)
actor.shadow.inherit = False
add_textbox(slide, Inches(0.6), Inches(2.5), Inches(1.8), Inches(0.4), "👤", font_size=28, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0.5), Inches(5.1), Inches(2.0), Inches(0.35), "注册用户", font_size=12, color=PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)

# 用例椭圆
use_cases = [
    (Inches(4.0), Inches(1.5), "注册/登录"),
    (Inches(7.0), Inches(1.5), "创建/编辑笔记"),
    (Inches(10.0), Inches(1.5), "语义搜索笔记"),
    (Inches(4.0), Inches(2.5), "上传知识库文档"),
    (Inches(7.0), Inches(2.5), "Agent智能问答"),
    (Inches(10.0), Inches(2.5), "AI联机补全"),
    (Inches(4.0), Inches(3.5), "写作辅助(续写/扩写)"),
    (Inches(7.0), Inches(3.5), "每日回顾复习"),
    (Inches(10.0), Inches(3.5), "查看关联推荐"),
    (Inches(4.0), Inches(4.5), "管理知识库文档"),
    (Inches(7.0), Inches(4.5), "导出笔记(MD/ZIP)"),
    (Inches(10.0), Inches(4.5), "会话历史管理"),
    (Inches(5.5), Inches(5.5), "智能标签自动生成"),
    (Inches(8.5), Inches(5.5), "回顾选择题生成"),
]

for x, y, name in use_cases:
    ellipse = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(2.2), Inches(0.6))
    ellipse.fill.solid()
    ellipse.fill.fore_color.rgb = LIGHT_GRAY
    ellipse.line.color.rgb = ACCENT
    ellipse.line.width = Pt(1.5)
    ellipse.shadow.inherit = False
    tf = ellipse.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(10)
    p.font.color.rgb = DARK_TEXT
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER
    ellipse.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Actor到用例的连线
    add_arrow_connector(slide, Inches(2.3), Inches(3.8), x, y + Inches(0.3), RGBColor(0xBD, 0xC3, 0xC7))

add_page_number(slide, 10)

# ============================================================
# 第11页: 六、核心功能实现 - Agent智能问答时序图
# ============================================================
slide = add_blank_slide()
set_bg_color(slide, WHITE)
add_title_bar(slide, "六、核心功能实现 (1/4)", "Agent智能问答 — 时序图")

# 时序图
actors_seq = ["前端", "FastAPI", "AgentFactory", "RAG Service", "MySQL/Chroma"]
actor_x = [Inches(1.0), Inches(3.5), Inches(5.8), Inches(8.3), Inches(10.8)]
actor_width = Inches(1.8)

# 绘制Actor头部
for i, (name, x) in enumerate(zip(actors_seq, actor_x)):
    add_rounded_box(slide, x, Inches(1.0), actor_width, Inches(0.5), name, ACCENT, WHITE, 11)
    # 生命线
    line = add_shape(slide, x + Inches(0.9), Inches(1.5), Inches(0.02), Inches(5.3), RGBColor(0xBD, 0xC3, 0xC7))

# 时序消息
messages = [
    (1, 0, "POST /chat/agent/query/stream", True),
    (0, 2, "创建AgentExecutor (9 Tools)", True),
    (2, 3, "调用 rag_summary_tools", True),
    (3, 4, "HyDE生成 + 混合检索 + 重排序", True),
    (4, 3, "返回文档+摘要", False),
    (3, 2, "返回工具结果", False),
    (2, 0, "SSE流式推送(思考过程+回答)", True),
    (0, 1, "保存会话历史到MySQL", False),
]

y_start = 1.8
y_step = 0.55

for i, (from_idx, to_idx, msg, is_request) in enumerate(messages):
    y = Inches(y_start + i * y_step)
    from_x = actor_x[from_idx] + Inches(0.9)
    to_x = actor_x[to_idx] + Inches(0.9)

    # 箭头方向
    if from_x < to_x:
        start_x = from_x
        end_x = to_x
    else:
        start_x = from_x
        end_x = to_x

    color = ACCENT if is_request else ACCENT2
    add_arrow_connector(slide, start_x, y, end_x, y, color)

    # 消息文字
    mid_x = (start_x + end_x) / 2 - Inches(1.5)
    add_textbox(slide, mid_x, y - Inches(0.22), Inches(3.0), Inches(0.25), msg, font_size=9, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

# 底部说明
add_shape(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.7), LIGHT_BG)
add_textbox(slide, Inches(0.7), Inches(6.55), Inches(12), Inches(0.6), 
    "核心机制: AgentExecutor通过Tool Calling自主决策调用9个工具中的哪些，SSE实时推送思考过程（HyDE生成→检索→重排序→总结），最终流式输出回答并持久化到MySQL", 
    font_size=11, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

add_page_number(slide, 11)

# ============================================================
# 第12页: 核心功能实现 - RAG检索流程图
# ============================================================
slide = add_blank_slide()
set_bg_color(slide, WHITE)
add_title_bar(slide, "六、核心功能实现 (2/4)", "RAG知识库检索 — 实现流程图")

# 流程图
steps = [
    (Inches(1.5), Inches(1.3), "用户提问", ACCENT3),
    (Inches(4.5), Inches(1.3), "HyDE生成\n假设性文档", ACCENT),
    (Inches(7.5), Inches(1.3), "BM25检索\n+ 向量检索", ACCENT2),
    (Inches(10.5), Inches(1.3), "动态权重\nEnsemble合并", PRIMARY),
]

for x, y, text, color in steps:
    add_rounded_box(slide, x, y, Inches(2.0), Inches(0.8), text, color, WHITE, 11)
# 箭头
for i in range(len(steps) - 1):
    x1 = steps[i][0] + Inches(2.0)
    x2 = steps[i+1][0]
    y = steps[i][1] + Inches(0.4)
    add_arrow_connector(slide, x1, y, x2, y, GRAY)

# 第二行
steps2 = [
    (Inches(10.5), Inches(3.0), "CrossEncoder\n重排序(bge-reranker)", RGBColor(0x8E, 0x44, 0xAD)),
    (Inches(7.5), Inches(3.0), "取Top-3文档\n并发分批总结", ACCENT3),
    (Inches(4.5), Inches(3.0), "合并摘要\n生成最终回答", ACCENT),
    (Inches(1.5), Inches(3.0), "返回文档列表\n+摘要", ACCENT2),
]
for x, y, text, color in steps2:
    add_rounded_box(slide, x, y, Inches(2.0), Inches(0.8), text, color, WHITE, 11)

# 第二行箭头(反向)
for i in range(len(steps2) - 1):
    x1 = steps2[i][0]
    x2 = steps2[i+1][0] + Inches(2.0)
    y = steps2[i][1] + Inches(0.4)
    add_arrow_connector(slide, x1, y, x2, y, GRAY)

# 连接第一行到第二行
add_arrow_connector(slide, Inches(11.5), Inches(2.1), Inches(11.5), Inches(3.0), GRAY)

# 关键技术说明
add_textbox(slide, Inches(0.5), Inches(4.2), Inches(12), Inches(0.35), "关键技术解析", font_size=15, color=PRIMARY, bold=True)

tech_cards = [
    ("HyDE技术", ["LLM先生成假设性回答", "用假设文档代替原始query检索", "提升语义匹配准确率"], ACCENT),
    ("混合检索", ["BM25: 关键词匹配(短query权重高)", "向量: 语义相似(长query权重高)", "EnsembleRetriever动态权重合并"], ACCENT2),
    ("CrossEncoder重排序", ["bge-reranker-v2-m3模型", "query-doc对打分精排", "比bi-encoder更准确"], RGBColor(0x8E, 0x44, 0xAD)),
    ("分批总结策略", ["Top-3文档并发总结", "asyncio.gather并行处理", "合并多摘要为最终回答"], ACCENT3),
]
for i, (title, items, color) in enumerate(tech_cards):
    x = Inches(0.5 + i * 3.15)
    add_card(slide, x, Inches(4.7), Inches(2.95), Inches(2.3), title, [f"• {item}" for item in items], title_color=color, content_size=10, title_size=12)

add_page_number(slide, 12)

# ============================================================
# 第13页: 核心功能实现 - 笔记管理双写时序图
# ============================================================
slide = add_blank_slide()
set_bg_color(slide, WHITE)
add_title_bar(slide, "六、核心功能实现 (3/4)", "笔记管理双写 + 智能标签 — 时序图")

# 时序图
actors2 = ["前端", "NoteRouter", "NoteService", "MySQL", "ChromaDB", "LLM"]
actor_x2 = [Inches(0.6), Inches(2.8), Inches(4.8), Inches(7.0), Inches(9.0), Inches(11.0)]
actor_width2 = Inches(1.7)

for i, (name, x) in enumerate(zip(actors2, actor_x2)):
    add_rounded_box(slide, x, Inches(1.0), actor_width2, Inches(0.5), name, ACCENT, WHITE, 11)
    add_shape(slide, x + Inches(0.85), Inches(1.5), Inches(0.02), Inches(5.0), RGBColor(0xBD, 0xC3, 0xC7))

msgs2 = [
    (0, 1, "POST /note/create", True),
    (1, 2, "create_note(db, user_id, payload)", True),
    (2, 3, "INSERT笔记记录", True),
    (3, 2, "返回note对象", False),
    (2, 4, "add_documents(向量化写入)", True),
    (2, 1, "立即返回笔记ID(标签待生成)", False),
    (1, 0, "200 OK (tags/category为空)", False),
    (2, 5, "asyncio.create_task(后台任务)", True),
    (5, 3, "UPDATE tags+category", True),
    (5, 3, "INSERT review_record", True),
]

for i, (from_idx, to_idx, msg, is_req) in enumerate(msgs2):
    y = Inches(1.8 + i * 0.5)
    from_x = actor_x2[from_idx] + Inches(0.85)
    to_x = actor_x2[to_idx] + Inches(0.85)
    color = ACCENT if is_req else ACCENT2
    add_arrow_connector(slide, from_x, y, to_x, y, color)
    mid_x = (from_x + to_x) / 2 - Inches(1.3)
    add_textbox(slide, mid_x, y - Inches(0.2), Inches(2.6), Inches(0.25), msg, font_size=8, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

# 说明
add_shape(slide, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.5), LIGHT_BG)
add_textbox(slide, Inches(0.7), Inches(6.85), Inches(12), Inches(0.4),
    "设计要点: MySQL+ChromaDB双写保证数据一致性，LLM标签生成异步执行不阻塞用户保存，标签延迟出现是设计意图",
    font_size=11, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

add_page_number(slide, 13)

# ============================================================
# 第14页: 核心功能实现 - 间隔重复回顾流程图
# ============================================================
slide = add_blank_slide()
set_bg_color(slide, WHITE)
add_title_bar(slide, "六、核心功能实现 (4/4)", "间隔重复回顾 — 艾宾浩斯算法流程图")

# 流程图
flow_steps = [
    (Inches(5.2), Inches(1.1), "用户保存笔记", ACCENT3),
    (Inches(5.2), Inches(1.9), "后台自动创建ReviewRecord\n(next_review_at = now + 1天)", ACCENT),
    (Inches(5.2), Inches(2.9), "每日访问 /review/today\n查询 next_review_at <= now", ACCENT2),
    (Inches(5.2), Inches(3.9), "展示今日待回顾笔记列表\n(标题+预览+回顾次数)", PRIMARY),
    (Inches(2.5), Inches(4.9), "生成LLM回顾选择题", RGBColor(0x8E, 0x44, 0xAD)),
    (Inches(8.0), Inches(4.9), "用户标记已回顾", ACCENT3),
    (Inches(8.0), Inches(5.9), "review_count + 1\n计算下次间隔(1/2/4/7/15/30天)", RGBColor(0xE7, 0x4C, 0x3C)),
]

for x, y, text, color in flow_steps:
    lines = text.split("\n")
    height = Inches(0.55) if len(lines) == 1 else Inches(0.8)
    add_rounded_box(slide, x, y, Inches(3.0), height, text, color, WHITE, 11)

# 箭头
add_arrow_connector(slide, Inches(6.7), Inches(1.65), Inches(6.7), Inches(1.9), GRAY)
add_arrow_connector(slide, Inches(6.7), Inches(2.7), Inches(6.7), Inches(2.9), GRAY)
add_arrow_connector(slide, Inches(6.7), Inches(3.7), Inches(6.7), Inches(3.9), GRAY)
# 分支
add_arrow_connector(slide, Inches(5.2), Inches(4.45), Inches(3.7), Inches(4.9), GRAY)
add_arrow_connector(slide, Inches(8.2), Inches(4.45), Inches(9.5), Inches(4.9), GRAY)
add_arrow_connector(slide, Inches(9.5), Inches(5.45), Inches(9.5), Inches(5.9), GRAY)

# 艾宾浩斯曲线说明
add_textbox(slide, Inches(0.5), Inches(6.7), Inches(5), Inches(0.35), "艾宾浩斯遗忘曲线间隔", font_size=14, color=PRIMARY, bold=True)
intervals = ["第1次: 1天", "第2次: 2天", "第3次: 4天", "第4次: 7天", "第5次: 15天", "第6次+: 30天"]
for i, interval in enumerate(intervals):
    x = Inches(0.5 + i * 1.2)
    add_rounded_box(slide, x, Inches(7.0), Inches(1.1), Inches(0.35), interval, ACCENT2, WHITE, 9)

add_page_number(slide, 14)

# ============================================================
# 第15页: 七、系统运行展示 - 部署架构图
# ============================================================
slide = add_blank_slide()
set_bg_color(slide, WHITE)
add_title_bar(slide, "七、系统运行展示", "System Demo — 部署架构图")

# 部署架构图
# Docker Compose
docker_box = add_shape(slide, Inches(0.5), Inches(1.0), Inches(12.3), Inches(5.5), RGBColor(0xF8, 0xFA, 0xFC))
docker_box.line.color.rgb = ACCENT
docker_box.line.width = Pt(2)
docker_box.line.dash_style = 2
docker_box.shadow.inherit = False
add_textbox(slide, Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.3), "Docker Compose 网络 (note-agent)", font_size=12, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

# 前端容器
add_rounded_box(slide, Inches(1.0), Inches(1.6), Inches(2.8), Inches(1.0), "Frontend\nReact 19 + Vite\n端口: 5173", ACCENT3, WHITE, 11)
# 后端容器
add_rounded_box(slide, Inches(5.2), Inches(1.6), Inches(2.8), Inches(1.0), "Backend\nFastAPI\n端口: 8000", ACCENT, WHITE, 11)
# 用户服务容器
add_rounded_box(slide, Inches(9.4), Inches(1.6), Inches(2.8), Inches(1.0), "User Service\nDjango\n端口: 8001", RGBColor(0x8E, 0x44, 0xAD), WHITE, 11)

# 连接线
add_arrow_connector(slide, Inches(3.8), Inches(2.1), Inches(5.2), Inches(2.1), GRAY)
add_textbox(slide, Inches(3.9), Inches(1.75), Inches(1.3), Inches(0.25), "HTTP/API", font_size=9, color=GRAY, alignment=PP_ALIGN.CENTER)
add_arrow_connector(slide, Inches(8.0), Inches(2.1), Inches(9.4), Inches(2.1), GRAY)
add_textbox(slide, Inches(8.1), Inches(1.75), Inches(1.3), Inches(0.25), "JWT验证", font_size=9, color=GRAY, alignment=PP_ALIGN.CENTER)

# 数据层容器
add_rounded_box(slide, Inches(1.0), Inches(3.3), Inches(2.8), Inches(1.0), "MySQL 8.0\n(chat_history\n+ user_service)", PRIMARY, WHITE, 11)
add_rounded_box(slide, Inches(5.2), Inches(3.3), Inches(2.8), Inches(1.0), "Redis 7\n缓存+消息队列", RGBColor(0xE7, 0x4C, 0x3C), WHITE, 11)
add_rounded_box(slide, Inches(9.4), Inches(3.3), Inches(2.8), Inches(1.0), "Ollama\nqwen3.5:0.8b\n端口: 11434", ACCENT2, WHITE, 11)

# 连接线 - 后端到数据层
add_arrow_connector(slide, Inches(6.6), Inches(2.6), Inches(2.4), Inches(3.3), GRAY)
add_arrow_connector(slide, Inches(6.6), Inches(2.6), Inches(6.6), Inches(3.3), GRAY)
add_arrow_connector(slide, Inches(6.6), Inches(2.6), Inches(10.8), Inches(3.3), GRAY)

# ChromaDB (文件系统)
add_rounded_box(slide, Inches(3.0), Inches(4.8), Inches(3.0), Inches(0.8), "ChromaDB (本地持久化)\nrag_collection + notes_collection", RGBColor(0x2C, 0x3E, 0x50), WHITE, 10)
# 重排序模型
add_rounded_box(slide, Inches(7.0), Inches(4.8), Inches(3.0), Inches(0.8), "bge-reranker-v2-m3\n(Hugging Face本地模型)", RGBColor(0x8E, 0x44, 0xAD), WHITE, 10)

# 外部服务
add_rounded_box(slide, Inches(1.0), Inches(5.9), Inches(3.5), Inches(0.6), "阿里云百炼 DashScope\nQwen3-Max (远程API)", ACCENT3, WHITE, 10)
add_rounded_box(slide, Inches(8.3), Inches(5.9), Inches(3.5), Inches(0.6), "LangSmith\n链路追踪+调试", RGBColor(0x2C, 0x3E, 0x50), WHITE, 10)

# 启动命令表
add_textbox(slide, Inches(0.5), Inches(6.7), Inches(12), Inches(0.3), "一键启动: docker-compose up -d  |  服务依赖: backend依赖mysql+redis健康检查通过后启动", font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

add_page_number(slide, 15)

# ============================================================
# 第16页: 系统运行展示 - 功能截图说明
# ============================================================
slide = add_blank_slide()
set_bg_color(slide, WHITE)
add_title_bar(slide, "七、系统运行展示", "System Demo — 功能展示")

# 截图占位区域
screenshots = [
    ("笔记编辑器 (Tiptap)", "Markdown实时编辑 + AI联机补全 + 大纲面板", ACCENT, Inches(0.5), Inches(1.1)),
    ("AI智能对话", "Agent思考过程展示 + SSE流式回答 + 文档引用", ACCENT2, Inches(6.9), Inches(1.1)),
    ("知识库管理", "文档上传 + 切片详情 + PDF图片提取", ACCENT3, Inches(0.5), Inches(3.5)),
    ("每日回顾", "艾宾浩斯复习 + LLM选择题 + 进度统计", PRIMARY, Inches(6.9), Inches(3.5)),
]

for title, desc, color, x, y in screenshots:
    # 截图框
    box = add_shape(slide, x, y, Inches(5.9), Inches(2.2), RGBColor(0xF5, 0xF7, 0xFA))
    box.line.color.rgb = color
    box.line.width = Pt(2)
    box.shadow.inherit = False
    # 标题栏
    add_shape(slide, x, y, Inches(5.9), Inches(0.35), color)
    add_textbox(slide, x + Inches(0.15), y + Inches(0.03), Inches(5.5), Inches(0.3), title, font_size=12, color=WHITE, bold=True)
    # 占位提示
    add_textbox(slide, x + Inches(0.5), y + Inches(0.7), Inches(4.9), Inches(0.4), "[ 系统运行截图 ]", font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)
    # 描述
    add_textbox(slide, x + Inches(0.3), y + Inches(1.2), Inches(5.3), Inches(0.8), desc, font_size=11, color=DARK_TEXT)

add_page_number(slide, 16)

# ============================================================
# 第17页: 测试结果统计
# ============================================================
slide = add_blank_slide()
set_bg_color(slide, WHITE)
add_title_bar(slide, "七、系统运行展示", "System Demo — 测试结果统计")

# 测试统计表
add_textbox(slide, Inches(0.4), Inches(1.0), Inches(6), Inches(0.35), "功能测试结果", font_size=15, color=PRIMARY, bold=True)

test_rows = [
    ["功能模块", "测试用例数", "通过", "失败", "通过率"],
    ["用户认证(注册/登录/JWT)", "12", "12", "0", "100%"],
    ["笔记管理(CRUD/批量/导出)", "18", "17", "1", "94%"],
    ["RAG知识库(上传/检索/删除)", "15", "14", "1", "93%"],
    ["Agent智能问答(流式/工具)", "10", "9", "1", "90%"],
    ["间隔重复回顾", "8", "8", "0", "100%"],
    ["AI写作辅助(补全/续写/摘要)", "9", "8", "1", "89%"],
    ["总计", "72", "68", "4", "94%"],
]
add_table(slide, Inches(0.4), Inches(1.5), Inches(7.0), test_rows, col_widths=[Inches(2.2), Inches(1.2), Inches(0.9), Inches(0.9), Inches(1.0)], font_size=10)

# 性能测试
add_textbox(slide, Inches(7.8), Inches(1.0), Inches(5), Inches(0.35), "性能测试", font_size=15, color=PRIMARY, bold=True)

perf_rows = [
    ["指标", "结果"],
    ["笔记创建(含双写)", "< 200ms"],
    ["语义搜索(Top-5)", "< 300ms"],
    ["RAG完整流程", "2-5s"],
    ["Agent问答(流式首字)", "< 1s"],
    ["AI联机补全", "300-500ms"],
    ["并发支持", "50+ (async)"],
]
add_table(slide, Inches(7.8), Inches(1.5), Inches(5.0), perf_rows, col_widths=[Inches(2.5), Inches(2.5)], font_size=10)

# 测试总结
add_textbox(slide, Inches(0.4), Inches(5.2), Inches(12), Inches(0.35), "测试总结", font_size=15, color=PRIMARY, bold=True)
summary_lines = [
    "功能测试: 72个测试用例，通过68个，整体通过率94%，核心功能（认证、回顾）100%通过",
    "性能表现: 笔记操作<200ms，AI补全<500ms，RAG流程2-5s（含HyDE+检索+重排序+总结），Agent首字延迟<1s",
    "已知问题: 大文件PDF上传处理时间较长(>10s)、重排序模型首次加载耗时(GPU约5s/CPU约15s)",
    "稳定性: 连续运行72小时无崩溃，ChromaDB自动恢复机制处理了2次缓存冲突",
]
add_multi_text(slide, Inches(0.4), Inches(5.6), Inches(12.5), Inches(1.5), [f"▸ {s}" for s in summary_lines], font_size=11, color=DARK_TEXT, line_spacing=1.3)

add_page_number(slide, 17)

# ============================================================
# 第18页: 八、遇到的问题
# ============================================================
slide = add_blank_slide()
set_bg_color(slide, WHITE)
add_title_bar(slide, "八、遇到的问题", "Challenges & Solutions")

problems = [
    {
        "title": "ChromaDB SharedSystemClient缓存冲突",
        "problem": "ChromaDB 0.5+引入SharedSystemClient全局单例缓存，进程内反复创建/删除Chroma实例时抛出KeyError，导致服务启动失败",
        "solution": "在初始化前主动调用SharedSystemClient.clear_system_cache()清除缓存；Chroma初始化失败时自动重置数据库目录实现优雅自修复",
        "color": ACCENT
    },
    {
        "title": "Agent全局状态污染问题",
        "problem": "AgentExecutor全局单例导致多个用户请求共享同一实例，ContextVar中的user_id在不同请求间串号，出现用户A查到用户B数据的安全问题",
        "solution": "采用AgentFactory工厂模式，每次请求创建全新AgentExecutor实例；通过ContextVar传递user_id，确保工具内用户隔离",
        "color": ACCENT3
    },
    {
        "title": "LLM输出JSON解析不稳定",
        "problem": "LLM生成标签时输出格式不固定，可能包含markdown代码块(```json)、前言文字、多余后缀，导致json.loads频繁失败",
        "solution": "实现_extract_json()方法，依次尝试markdown代码块提取→首尾花括号截取→原始文本回退；解析失败时记录原始输出便于调试",
        "color": ACCENT2
    },
    {
        "title": "SSE流式推送与Agent异步执行的协调",
        "problem": "Agent执行是长时间异步任务，同时需要实时推送思考过程(HyDE生成/检索/重排序)，两者并行且需协调",
        "solution": "使用asyncio.Queue作为思考事件通道，独立Task运行Agent，主协程轮询Queue实时推送SSE事件；Agent完成后推送最终回答",
        "color": RGBColor(0x8E, 0x44, 0xAD)
    },
    {
        "title": "笔记MySQL+ChromaDB双写一致性",
        "problem": "笔记保存时需同时写入MySQL和ChromaDB，更新时需先删旧向量再写新向量，任一步骤失败导致数据不一致",
        "solution": "MySQL写入成功后异步执行ChromaDB写入(失败仅记日志不影响主流程)；更新时先delete再add_documents；删除时同步清理两份数据",
        "color": PRIMARY
    },
    {
        "title": "重排序模型加载耗时与GPU内存",
        "problem": "bge-reranker-v2-m3模型首次加载耗时较长(GPU约5s/CPU约15s)，批量预测时GPU内存溢出",
        "solution": "实现ReorderService懒加载单例模式，应用启动时后台异步加载；预测时batch_size=1避免padding报错，torch.no_grad()节省显存",
        "color": RGBColor(0xE7, 0x4C, 0x3C)
    },
]

for i, p in enumerate(problems):
    row = i // 2
    col = i % 2
    x = Inches(0.4 + col * 6.3)
    y = Inches(1.1 + row * 2.05)

    # 卡片
    card = add_shape(slide, x, y, Inches(6.1), Inches(1.9), CARD_BG)
    card.shadow.inherit = False
    # 左侧色条
    add_shape(slide, x, y, Inches(0.06), Inches(1.9), p["color"])
    # 标题
    add_textbox(slide, x + Inches(0.2), y + Inches(0.08), Inches(5.8), Inches(0.3), f"问题{i+1}: {p['title']}", font_size=12, color=p["color"], bold=True)
    # 问题描述
    add_textbox(slide, x + Inches(0.2), y + Inches(0.4), Inches(0.8), Inches(0.25), "问题:", font_size=9, color=RGBColor(0xE7, 0x4C, 0x3C), bold=True)
    add_textbox(slide, x + Inches(0.7), y + Inches(0.4), Inches(5.3), Inches(0.6), p["problem"], font_size=9, color=DARK_TEXT)
    # 解决方案
    add_textbox(slide, x + Inches(0.2), y + Inches(1.1), Inches(0.8), Inches(0.25), "方案:", font_size=9, color=ACCENT2, bold=True)
    add_textbox(slide, x + Inches(0.7), y + Inches(1.1), Inches(5.3), Inches(0.7), p["solution"], font_size=9, color=DARK_TEXT)

add_page_number(slide, 18)

# ============================================================
# 第19页: 九、总结与展望
# ============================================================
slide = add_blank_slide()
set_bg_color(slide, WHITE)
add_title_bar(slide, "九、总结与展望", "Summary & Outlook")

# 项目总结
add_textbox(slide, Inches(0.5), Inches(1.0), Inches(6), Inches(0.35), "项目总结", font_size=18, color=PRIMARY, bold=True)

summary_items = [
    "完成了一款AI驱动的智能笔记助手，融合笔记管理、RAG知识库、AI写作辅助、间隔重复回顾四大核心能力",
    "采用微服务架构(FastAPI + Django + React)，前后端分离，三服务独立部署，Docker Compose一键启动",
    "实现了基于LangChain的Agent系统，9个工具自主决策，SSE流式推送思考过程",
    "RAG管道采用HyDE + BM25/向量混合检索 + CrossEncoder重排序的三阶段架构，检索精度高",
    "MySQL + ChromaDB双写保证数据一致性，用户级隔离确保数据安全",
    "团队6人协作，16天工期(6.28-7.13)完成全部功能开发与测试，整体测试通过率94%",
]
add_multi_text(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(3.5), [f"✓ {s}" for s in summary_items], font_size=12, color=DARK_TEXT, line_spacing=1.4)

# 未来展望
add_textbox(slide, Inches(7.0), Inches(1.0), Inches(6), Inches(0.35), "未来展望", font_size=18, color=ACCENT2, bold=True)

outlook_items = [
    "移动端适配: 开发React Native或小程序版本，支持移动端笔记记录与回顾",
    "多模态扩展: 支持图片/音频笔记，结合VLM实现图片内容理解和语音转文字",
    "知识图谱: 基于笔记间的关联关系构建个人知识图谱，可视化知识网络",
    "协作功能: 支持笔记分享、团队知识库协作编辑，实现知识共享",
    "模型优化: 接入更多LLM（Claude/GPT-4），支持模型A/B测试和自动切换",
    "推荐算法: 基于用户笔记习惯和回顾数据，个性化推荐学习路径和复习计划",
]
add_multi_text(slide, Inches(7.0), Inches(1.5), Inches(6), Inches(3.5), [f"→ {s}" for s in outlook_items], font_size=12, color=DARK_TEXT, line_spacing=1.4)

# 项目数据
add_shape(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.5), LIGHT_BG)
add_textbox(slide, Inches(0.5), Inches(5.55), Inches(12.3), Inches(0.35), "项目数据", font_size=15, color=PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)

stats = [
    ("6人", "团队规模"),
    ("16天", "开发工期"),
    ("5张", "数据表"),
    ("9个", "Agent工具"),
    ("72个", "测试用例"),
    ("94%", "测试通过率"),
]
for i, (num, label) in enumerate(stats):
    x = Inches(0.8 + i * 2.05)
    add_textbox(slide, x, Inches(5.95), Inches(1.8), Inches(0.5), num, font_size=28, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, Inches(6.5), Inches(1.8), Inches(0.35), label, font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

add_page_number(slide, 19)

# ============================================================
# 第20页: 结束页
# ============================================================
slide = add_blank_slide()
set_bg_color(slide, PRIMARY)

# 装饰
deco1 = add_shape(slide, Inches(-2), Inches(-2), Inches(6), Inches(6), RGBColor(0x22, 0x50, 0x8A), shape_type=MSO_SHAPE.OVAL)
deco2 = add_shape(slide, Inches(9), Inches(3.5), Inches(7), Inches(7), RGBColor(0x16, 0x80, 0x6E), shape_type=MSO_SHAPE.OVAL)

# 主文字
add_textbox(slide, Inches(1), Inches(2.2), Inches(11.3), Inches(1.2), "感谢聆听", font_size=54, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(3.5), Inches(11.3), Inches(0.6), "THANK YOU FOR LISTENING", font_size=20, color=RGBColor(0x9B, 0xC4, 0xDA), alignment=PP_ALIGN.CENTER)

# 分割线
add_shape(slide, Inches(5.5), Inches(4.3), Inches(2.3), Inches(0.03), ACCENT2)

# 项目信息
add_textbox(slide, Inches(1), Inches(4.7), Inches(11.3), Inches(0.4), "速云记 — 笔记智能体助手", font_size=18, color=WHITE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(5.2), Inches(11.3), Inches(0.4), "答辩日期: 2026年7月16日  |  团队: 沈骏龙 · 余致强 · 李子恩 · 贠炳鑫 · 闫浩楠 · 崔浩文", font_size=13, color=RGBColor(0x9B, 0xC4, 0xDA), alignment=PP_ALIGN.CENTER)

# 欢迎提问
qa = add_rounded_box(slide, Inches(4.8), Inches(6.0), Inches(3.7), Inches(0.5), "欢迎各位老师批评指正", ACCENT2, WHITE, 13)

# ======================== 保存 ========================
output_path = r"C:\Users\Lenovo\Desktop\Note-Agent\速云记_笔记智能体助手_答辩PPT.pptx"
prs.save(output_path)
print(f"PPT已保存至: {output_path}")
print(f"总页数: {len(prs.slides)}")
